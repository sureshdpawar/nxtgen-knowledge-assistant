from pathlib import Path

from pptx import (
    Presentation,
)

from app.parsers.base import (
    BaseParser,
    ParsedResult,
)


class PptxParser(
    BaseParser
):

    def extract(
        self,
        file_path: Path,
    ) -> ParsedResult:

        presentation = (
            Presentation(
                str(
                    file_path
                )
            )
        )

        pages: list[dict] = []

        for (
            slide_index,
            slide,
        ) in enumerate(
            presentation.slides,
            start=1,
        ):

            text_parts: list[str] = []

            for shape in slide.shapes:

                if (
                    getattr(
                        shape,
                        "has_text_frame",
                        False,
                    )
                ):
                    text = (
                        shape.text
                        .strip()
                    )

                    if text:
                        text_parts.append(
                            text
                        )

                if (
                    getattr(
                        shape,
                        "has_table",
                        False,
                    )
                ):
                    table = (
                        shape.table
                    )

                    for row in table.rows:

                        values = [
                            cell.text.strip()
                            for cell
                            in row.cells
                        ]

                        if any(
                            values
                        ):
                            text_parts.append(
                                " | ".join(
                                    values
                                )
                            )

            #
            # Include speaker notes when
            # available.
            #
            if (
                slide.has_notes_slide
            ):
                notes_frame = (
                    slide
                    .notes_slide
                    .notes_text_frame
                )

                if notes_frame:
                    notes = (
                        notes_frame.text
                        .strip()
                    )

                    if notes:
                        text_parts.append(
                            "Speaker Notes:\n"
                            + notes
                        )

            text = "\n".join(
                text_parts
            ).strip()

            pages.append(
                {
                    "page":
                        slide_index,

                    "text":
                        text,
                }
            )

        return {
            "pages":
                pages,

            "metadata": {
                "page_count":
                    len(
                        presentation
                        .slides
                    ),

                "document_type":
                    "presentation",
            },
        }